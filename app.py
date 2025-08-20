import os
import time
import uuid
from io import BytesIO
from zipfile import ZipFile, ZIP_DEFLATED
from typing import List, Tuple
from flask import Flask, request, jsonify, render_template, send_from_directory, url_for, abort
from werkzeug.utils import secure_filename
from gtts import gTTS
from mutagen.mp3 import MP3
from PIL import Image
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from openpyxl import Workbook
from moviepy.video.io.VideoFileClip import VideoFileClip
from rembg import remove as rembg_remove, new_session as rembg_new_session

# Flask application setup
app = Flask(__name__)

BASE_DIR = os.path.abspath(os.path.dirname(__file__))
AUDIO_DIR = os.path.join(BASE_DIR, 'generated')
UPLOAD_DIR = os.path.join(BASE_DIR, 'uploads')
os.makedirs(AUDIO_DIR, exist_ok=True)
os.makedirs(UPLOAD_DIR, exist_ok=True)

# Configuration
MAX_TEXT_LENGTH = int(os.environ.get('MAX_TEXT_LENGTH', 5000))
CLEANUP_MAX_AGE_SECONDS = int(os.environ.get('CLEANUP_MAX_AGE_SECONDS', 3600))  # 1 hour by default
GTTS_TLD = os.environ.get('GTTS_TLD', 'com')
MIN_DURATION_SECONDS = float(os.environ.get('MIN_DURATION_SECONDS', 0.3))

# Cleanup extensions handled
CLEANUP_EXTENSIONS = {
	'.mp3', '.wav', '.ogg', '.pdf', '.png', '.jpg', '.jpeg', '.zip', '.docx', '.pptx', '.xlsx'
}

# Preload lightweight rembg session for speed (u2netp is small & fast)
REMBG_MODEL = os.environ.get('REMBG_MODEL', 'u2netp')
REMBG_SESSION = rembg_new_session(REMBG_MODEL)


def clean_old_files() -> None:
	"""Delete generated/uploaded files older than CLEANUP_MAX_AGE_SECONDS."""
	try:
		def _clean_dir(target_dir: str) -> None:
			now = time.time()
			for entry in os.scandir(target_dir):
				if not entry.is_file():
					continue
				_, ext = os.path.splitext(entry.name)
				if ext.lower() not in CLEANUP_EXTENSIONS:
					continue
				try:
					mtime = entry.stat().st_mtime
					if now - mtime > CLEANUP_MAX_AGE_SECONDS:
						os.remove(entry.path)
				except FileNotFoundError:
					continue
		for d in (AUDIO_DIR, UPLOAD_DIR):
			_clean_dir(d)
	except Exception:
		pass


def generate_unique_filename(suffix: str = '') -> str:
	return f"tool_{uuid.uuid4().hex}{suffix}"


def parse_pages(pages: str, total: int) -> List[int]:
	"""Parse page selectors like '1-3,5' into 0-based page indices within total."""
	if not pages:
		return list(range(total))
	result = []
	for part in pages.split(','):
		part = part.strip()
		if not part:
			continue
		if '-' in part:
			a, b = part.split('-', 1)
			try:
				start = max(1, int(a))
				end = min(total, int(b))
			except ValueError:
				continue
			for p in range(start, end + 1):
				result.append(p - 1)
		else:
			try:
				p = int(part)
			except ValueError:
				continue
			if 1 <= p <= total:
				result.append(p - 1)
	# Deduplicate and sort
	return sorted(set(result))


@app.get('/')
def index():
	clean_old_files()
	return render_template('index.html')


@app.get('/tools')
def tools():
	clean_old_files()
	return render_template('tools.html')


# ========== Original TTS API ==========
@app.post('/api/convert')
def convert_text_to_mp3():
	clean_old_files()
	# Support JSON or form-encoded
	text = ''
	lang = 'en'
	if request.is_json:
		payload = request.get_json(silent=True) or {}
		text = (payload.get('text') or '').strip()
		lang = (payload.get('lang') or 'en').strip() or 'en'
	else:
		text = (request.form.get('text') or '').strip()
		lang = (request.form.get('lang') or 'en').strip() or 'en'

	if not text:
		return jsonify({'success': False, 'error': 'Text is required.'}), 400
	if len(text) > MAX_TEXT_LENGTH:
		return jsonify({'success': False, 'error': f'Text exceeds maximum length of {MAX_TEXT_LENGTH} characters.'}), 413

	filename = generate_unique_filename('.mp3')
	file_path = os.path.join(AUDIO_DIR, filename)

	try:
		tts = gTTS(text=text, lang=lang, tld=GTTS_TLD)
		tts.save(file_path)
		# Validate duration
		try:
			audio = MP3(file_path)
			duration_seconds = float(getattr(audio.info, 'length', 0.0) or 0.0)
			if duration_seconds < MIN_DURATION_SECONDS:
				os.remove(file_path)
				return jsonify({'success': False, 'error': 'Generated audio appears to be empty. Please try again.'}), 502
		except Exception:
			if not os.path.exists(file_path) or os.path.getsize(file_path) == 0:
				return jsonify({'success': False, 'error': 'Failed to validate generated audio.'}), 502
	except Exception as exc:
		try:
			if os.path.exists(file_path):
				os.remove(file_path)
		except Exception:
			pass
		return jsonify({'success': False, 'error': f'Failed to generate audio: {str(exc)}'}), 500

	download_url = url_for('download_file', filename=filename)
	stream_url = url_for('stream_file', filename=filename)
	return jsonify({'success': True, 'filename': filename, 'url': download_url, 'stream_url': stream_url})


# ========== Images → PDF ==========
@app.get('/tools/images-to-pdf')
def images_to_pdf_page():
	clean_old_files()
	return render_template('tools/images_to_pdf.html')


@app.post('/api/images-to-pdf')
def images_to_pdf_api():
	clean_old_files()
	files = request.files.getlist('images')
	page_size = (request.form.get('page_size') or 'A4').upper()
	order = request.form.get('order') or ''  # comma-separated original indices (1-based)
	if not files:
		return jsonify({'success': False, 'error': 'No images uploaded.'}), 400

	# Reorder files if order specified
	if order:
		try:
			indices = [int(x.strip()) - 1 for x in order.split(',') if x.strip()]
			files = [files[i] for i in indices if 0 <= i < len(files)]
		except Exception:
			pass

	size_map = {
		'A4': (2480, 3508),  # 300 DPI
		'LETTER': (2550, 3300),
	}
	page_px = size_map.get(page_size, size_map['A4'])

	pil_images: List[Image.Image] = []
	for f in files:
		try:
			img = Image.open(f.stream).convert('RGB')
			# Fit to page
			img.thumbnail(page_px, Image.LANCZOS)
			canvas = Image.new('RGB', page_px, 'white')
			offset = ((page_px[0] - img.size[0]) // 2, (page_px[1] - img.size[1]) // 2)
			canvas.paste(img, offset)
			pil_images.append(canvas)
		except Exception:
			continue

	if not pil_images:
		return jsonify({'success': False, 'error': 'Failed to process images.'}), 400

	out_name = generate_unique_filename('.pdf')
	out_path = os.path.join(AUDIO_DIR, out_name)
	pil_images[0].save(out_path, save_all=True, append_images=pil_images[1:])
	return jsonify({'success': True, 'url': url_for('download_file', filename=out_name), 'filename': out_name})


# ========== PDF → Images ==========
@app.get('/tools/pdf-to-images')

def pdf_to_images_page():
	clean_old_files()
	return render_template('tools/pdf_to_images.html')


@app.post('/api/pdf-to-images')
def pdf_to_images_api():
	clean_old_files()
	pdf_file = request.files.get('pdf')
	pages = request.form.get('pages') or ''
	quality = (request.form.get('quality') or 'medium').lower()  # low/medium/high
	if not pdf_file:
		return jsonify({'success': False, 'error': 'No PDF uploaded.'}), 400

	tmp_pdf = os.path.join(UPLOAD_DIR, secure_filename(generate_unique_filename('.pdf')))
	pdf_file.save(tmp_pdf)

	zoom_map = {'low': 1.0, 'medium': 2.0, 'high': 3.0}
	zoom = zoom_map.get(quality, 2.0)

	doc = fitz.open(tmp_pdf)
	page_indices = parse_pages(pages, doc.page_count)
	out_files: List[Tuple[int, str]] = []
	for i in page_indices:
		page = doc.load_page(i)
		mat = fitz.Matrix(zoom, zoom)
		pix = page.get_pixmap(matrix=mat, alpha=False)
		out_name = generate_unique_filename(f'_p{i+1}.png')
		out_path = os.path.join(AUDIO_DIR, out_name)
		pix.save(out_path)
		out_files.append((i + 1, out_name))
	doc.close()

	# Zip results
	zip_name = generate_unique_filename('.zip')
	zip_path = os.path.join(AUDIO_DIR, zip_name)
	with ZipFile(zip_path, 'w', ZIP_DEFLATED) as z:
		for _, fname in out_files:
			z.write(os.path.join(AUDIO_DIR, fname), arcname=fname)

	return jsonify({'success': True, 'zip_url': url_for('download_file', filename=zip_name)})


# ========== PDF → Office ==========
@app.get('/tools/pdf-to-office')

def pdf_to_office_page():
	clean_old_files()
	return render_template('tools/pdf_to_office.html')


@app.post('/api/pdf-to-office')
def pdf_to_office_api():
	clean_old_files()
	pdf_file = request.files.get('pdf')
	target = (request.form.get('target') or 'docx').lower()  # docx/xlsx/pptx
	pages = request.form.get('pages') or ''
	if not pdf_file:
		return jsonify({'success': False, 'error': 'No PDF uploaded.'}), 400

	tmp_pdf = os.path.join(UPLOAD_DIR, secure_filename(generate_unique_filename('.pdf')))
	pdf_file.save(tmp_pdf)
	doc = fitz.open(tmp_pdf)
	page_indices = parse_pages(pages, doc.page_count)

	if target == 'docx':
		out_name = generate_unique_filename('.docx')
		out_path = os.path.join(AUDIO_DIR, out_name)
		word = Document()
		for idx in page_indices:
			text = doc.load_page(idx).get_text()
			if word.paragraphs:
				word.add_page_break()
			word.add_paragraph(text)
		word.save(out_path)
		result_url = url_for('download_file', filename=out_name)
	elif target == 'pptx':
		out_name = generate_unique_filename('.pptx')
		out_path = os.path.join(AUDIO_DIR, out_name)
		prs = Presentation()
		# Standard 10"x7.5" size. We'll fit images to slide
		for idx in page_indices:
			page = doc.load_page(idx)
			pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
			img_bytes = pix.tobytes('png')
			slide = prs.slides.add_slide(prs.slide_layouts[6])
			# Save image to stream for python-pptx
			bio = BytesIO(img_bytes)
			pic = slide.shapes.add_picture(bio, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
		prs.save(out_path)
		result_url = url_for('download_file', filename=out_name)
	elif target == 'xlsx':
		out_name = generate_unique_filename('.xlsx')
		out_path = os.path.join(AUDIO_DIR, out_name)
		xl = Workbook()
		ws = xl.active
		ws.title = 'Page 1' if page_indices else 'Pages'
		for i, idx in enumerate(page_indices, start=1):
			ws = xl.create_sheet(title=f'Page {i}') if i > 1 else ws
			text = doc.load_page(idx).get_text()
			ws['A1'] = text
		xl.save(out_path)
		result_url = url_for('download_file', filename=out_name)
	else:
		doc.close()
		return jsonify({'success': False, 'error': 'Unsupported target format.'}), 400

	doc.close()
	return jsonify({'success': True, 'url': result_url})


# ========== Office → MP3 ==========
@app.get('/tools/office-to-mp3')

def office_to_mp3_page():
	clean_old_files()
	return render_template('tools/office_to_mp3.html')


def extract_text_from_docx(path: str) -> str:
	word = Document(path)
	return '\n'.join(p.text for p in word.paragraphs if p.text)


def extract_text_from_pptx(path: str) -> str:
	prs = Presentation(path)
	lines: List[str] = []
	for slide in prs.slides:
		for shape in slide.shapes:
			if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
				for p in shape.text_frame.paragraphs:
					lines.append(' '.join(r.text for r in p.runs))
	return '\n'.join([ln for ln in lines if ln.strip()])


@app.post('/api/office-to-mp3')
def office_to_mp3_api():
	clean_old_files()
	office = request.files.get('file')
	lang = (request.form.get('lang') or 'ar').strip() or 'ar'
	slow = (request.form.get('speed') or 'normal') == 'slow'
	tld = (request.form.get('tld') or GTTS_TLD).strip() or GTTS_TLD
	if not office:
		return jsonify({'success': False, 'error': 'No file uploaded.'}), 400

	# Preserve original extension even if the filename contains non-ASCII characters
	original_name = office.filename or ''
	orig_ext = os.path.splitext(original_name)[1].lower()
	if orig_ext not in ('.docx', '.pptx'):
		# Fallback: try inference from mimetype
		mt = (office.mimetype or '').lower()
		if 'word' in mt:
			orig_ext = '.docx'
		elif 'presentation' in mt or 'powerpoint' in mt:
			orig_ext = '.pptx'
		else:
			return jsonify({'success': False, 'error': 'Unsupported file type. Use .docx or .pptx'}), 400

	upload_name = generate_unique_filename(orig_ext)
	upload_path = os.path.join(UPLOAD_DIR, upload_name)
	office.save(upload_path)

	text = ''
	try:
		if orig_ext == '.docx':
			text = extract_text_from_docx(upload_path)
		elif orig_ext == '.pptx':
			text = extract_text_from_pptx(upload_path)
		if not text.strip():
			return jsonify({'success': False, 'error': 'Document appears to be empty.'}), 400

		out_name = generate_unique_filename('.mp3')
		out_path = os.path.join(AUDIO_DIR, out_name)
		tts = gTTS(text=text, lang=lang, tld=tld, slow=slow)
		tts.save(out_path)

		# Validate
		try:
			dur = float(getattr(MP3(out_path).info, 'length', 0.0) or 0.0)
			if dur < MIN_DURATION_SECONDS:
				os.remove(out_path)
				return jsonify({'success': False, 'error': 'Generated audio seems empty.'}), 502
		except Exception:
			if not os.path.exists(out_path) or os.path.getsize(out_path) == 0:
				return jsonify({'success': False, 'error': 'Failed to validate generated audio.'}), 502
	except Exception as exc:
		return jsonify({'success': False, 'error': f'Conversion failed: {str(exc)}'}), 500

	return jsonify({'success': True, 'url': url_for('download_file', filename=out_name), 'stream_url': url_for('stream_file', filename=out_name)})


# ========== Video → Audio ==========
@app.get('/tools/video-to-audio')

def video_to_audio_page():
	clean_old_files()
	return render_template('tools/video_to_audio.html')


@app.post('/api/video-to-audio')
def video_to_audio_api():
	clean_old_files()
	video = request.files.get('video')
	fmt = (request.form.get('format') or 'mp3').lower()  # mp3/wav/ogg
	bitrate = request.form.get('bitrate') or '192k'
	start = request.form.get('start') or ''
	end = request.form.get('end') or ''
	if not video:
		return jsonify({'success': False, 'error': 'No video uploaded.'}), 400

	# Preserve original extension for saved upload
	original_name = video.filename or ''
	orig_ext = os.path.splitext(original_name)[1].lower() or '.bin'
	upload_name = generate_unique_filename(orig_ext)
	upload_path = os.path.join(UPLOAD_DIR, upload_name)
	video.save(upload_path)

	try:
		clip = VideoFileClip(upload_path)
		if start or end:
			try:
				s = float(start) if start else 0.0
				e = float(end) if end else clip.duration
				if e > s:
					clip = clip.subclip(s, e)
			except Exception:
				pass
		out_ext = '.' + fmt
		out_name = generate_unique_filename(out_ext)
		out_path = os.path.join(AUDIO_DIR, out_name)
		if fmt == 'mp3':
			clip.audio.write_audiofile(out_path, bitrate=bitrate, codec='libmp3lame')
		elif fmt == 'wav':
			clip.audio.write_audiofile(out_path, codec='pcm_s16le')
		elif fmt == 'ogg':
			clip.audio.write_audiofile(out_path, bitrate=bitrate, codec='libvorbis')
		else:
			return jsonify({'success': False, 'error': 'Unsupported audio format.'}), 400
		clip.close()
	except Exception as exc:
		return jsonify({'success': False, 'error': f'Video processing failed (is ffmpeg installed?): {str(exc)}'}), 500

	return jsonify({'success': True, 'url': url_for('download_file', filename=out_name)})


# ========== Image Background Remover ==========
@app.get('/tools/background-remover')

def background_remover_page():
	clean_old_files()
	return render_template('tools/bg_remover.html')


@app.post('/api/background-remover')
def background_remover_api():
	clean_old_files()
	image = request.files.get('image')
	if not image:
		return jsonify({'success': False, 'error': 'No image uploaded.'}), 400
	try:
		inp = image.read()
		out = rembg_remove(inp, session=REMBG_SESSION)
		out_name = generate_unique_filename('.png')
		out_path = os.path.join(AUDIO_DIR, out_name)
		with open(out_path, 'wb') as f:
			f.write(out)
		return jsonify({'success': True, 'url': url_for('download_file', filename=out_name), 'filename': out_name, 'method': 'rembg', 'model': REMBG_MODEL})
	except Exception as exc:
		return jsonify({'success': False, 'error': f'Background removal failed: {str(exc)}'}), 500


# ========== Download/Stream and Errors ==========
@app.get('/download/<path:filename>')
def download_file(filename: str):
	basename = os.path.basename(filename)
	if basename != filename:
		abort(400)
	file_path = os.path.join(AUDIO_DIR, basename)
	if not os.path.isfile(file_path):
		abort(404)
	mtype_map = {
		'.mp3': 'audio/mpeg', '.wav': 'audio/wav', '.ogg': 'audio/ogg', '.pdf': 'application/pdf',
		'.png': 'image/png', '.jpg': 'image/jpeg', '.jpeg': 'image/jpeg', '.zip': 'application/zip',
		'.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
		'.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
		'.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
	}
	_, ext = os.path.splitext(basename)
	mtype = mtype_map.get(ext.lower(), 'application/octet-stream')
	return send_from_directory(AUDIO_DIR, basename, as_attachment=True, download_name=basename, mimetype=mtype)


@app.get('/stream/<path:filename>')
def stream_file(filename: str):
	basename = os.path.basename(filename)
	if basename != filename:
		abort(400)
	file_path = os.path.join(AUDIO_DIR, basename)
	if not os.path.isfile(file_path):
		abort(404)
	_, ext = os.path.splitext(basename)
	mtype = 'audio/mpeg' if ext.lower() == '.mp3' else 'application/octet-stream'
	return send_from_directory(AUDIO_DIR, basename, as_attachment=False, mimetype=mtype)


@app.errorhandler(404)
def not_found(_):
	return jsonify({'success': False, 'error': 'Not found'}), 404


@app.errorhandler(500)
def internal_error(_):
	return jsonify({'success': False, 'error': 'Internal server error'}), 500


if __name__ == '__main__':
	# For local development only; PythonAnywhere will use WSGI
	app.run(host='0.0.0.0', port=5000, debug=True) 